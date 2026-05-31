"""
Document parsing for a broad range of file formats.

Supported formats
-----------------
Plain text  : .txt, .md, .rst, .log
PDF         : .pdf  (requires pypdf)
Word        : .docx (requires python-docx)
Excel       : .xlsx, .xls (requires openpyxl / xlrd)
PowerPoint  : .pptx (requires python-pptx)
Spreadsheet : .csv, .tsv
Data        : .json, .jsonl
Web         : .html, .htm (requires beautifulsoup4)
Code/config : .py, .js, .ts, .yaml, .yml, .toml, .xml (plain-text fallback)

Unknown extensions fall back to plain-text reading.
"""

import csv
import json
import logging
from pathlib import Path
from typing import Any, Dict, Tuple

logger = logging.getLogger(__name__)

# All extensions that the pipeline will consider for ingestion
SUPPORTED_EXTENSIONS = {
    ".txt", ".md", ".rst", ".log",
    ".pdf",
    ".docx",
    ".xlsx", ".xls",
    ".pptx",
    ".csv", ".tsv",
    ".json", ".jsonl",
    ".html", ".htm",
    ".py", ".js", ".ts", ".yaml", ".yml", ".toml", ".xml",
}


def parse_document(file_path: str) -> Tuple[str, Dict[str, Any]]:
    """
    Parse a document file and return ``(text_content, metadata)``.

    Always returns a 2-tuple; ``text_content`` is an empty string on failure.
    """
    path = Path(file_path)
    suffix = path.suffix.lower()

    dispatch = {
        ".pdf":   _parse_pdf,
        ".docx":  _parse_docx,
        ".xlsx":  _parse_excel,
        ".xls":   _parse_excel,
        ".pptx":  _parse_pptx,
        ".csv":   _parse_csv,
        ".tsv":   _parse_tsv,
        ".json":  _parse_json,
        ".jsonl": _parse_jsonl,
        ".html":  _parse_html,
        ".htm":   _parse_html,
    }

    parser = dispatch.get(suffix, _parse_text)
    content, metadata = parser(file_path)
    metadata["file_name"] = path.name
    metadata["file_extension"] = suffix
    return content, metadata


# ── Format-specific parsers ───────────────────────────────────────────────────

def _parse_text(file_path: str) -> Tuple[str, Dict[str, Any]]:
    with open(file_path, encoding="utf-8", errors="ignore") as fh:
        content = fh.read()
    return content, {"format": "text"}


def _parse_pdf(file_path: str) -> Tuple[str, Dict[str, Any]]:
    try:
        import pypdf
        reader = pypdf.PdfReader(file_path)
        pages = [page.extract_text() or "" for page in reader.pages]
        content = "\n\n".join(p for p in pages if p.strip())
        return content, {"format": "pdf", "page_count": len(reader.pages)}
    except ImportError:
        logger.warning("pypdf not installed — falling back to plain-text for %s", file_path)
        return _parse_text(file_path)
    except Exception as exc:
        logger.error("PDF parsing failed for %s: %s", file_path, exc)
        return "", {"format": "pdf", "error": str(exc)}


def _parse_docx(file_path: str) -> Tuple[str, Dict[str, Any]]:
    try:
        import docx
        doc = docx.Document(file_path)
        # Paragraphs + tables
        lines = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                lines.append("\t".join(c.text for c in row.cells))
        content = "\n\n".join(lines)
        return content, {"format": "docx", "paragraph_count": len(doc.paragraphs)}
    except ImportError:
        logger.warning("python-docx not installed — falling back to plain-text for %s", file_path)
        return _parse_text(file_path)
    except Exception as exc:
        logger.error("DOCX parsing failed for %s: %s", file_path, exc)
        return "", {"format": "docx", "error": str(exc)}


def _parse_excel(file_path: str) -> Tuple[str, Dict[str, Any]]:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        lines = []
        sheet_count = 0
        for sheet in wb.worksheets:
            sheet_count += 1
            lines.append(f"## Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join("" if v is None else str(v) for v in row)
                if row_text.strip():
                    lines.append(row_text)
        content = "\n".join(lines)
        return content, {"format": "xlsx", "sheet_count": sheet_count}
    except ImportError:
        logger.warning("openpyxl not installed — falling back to plain-text for %s", file_path)
        return _parse_text(file_path)
    except Exception as exc:
        logger.error("Excel parsing failed for %s: %s", file_path, exc)
        return "", {"format": "xlsx", "error": str(exc)}


def _parse_pptx(file_path: str) -> Tuple[str, Dict[str, Any]]:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [
                shape.text for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if texts:
                slides.append(f"## Slide {i}\n" + "\n".join(texts))
        content = "\n\n".join(slides)
        return content, {"format": "pptx", "slide_count": len(prs.slides)}
    except ImportError:
        logger.warning("python-pptx not installed — falling back to plain-text for %s", file_path)
        return _parse_text(file_path)
    except Exception as exc:
        logger.error("PPTX parsing failed for %s: %s", file_path, exc)
        return "", {"format": "pptx", "error": str(exc)}


def _parse_csv(file_path: str) -> Tuple[str, Dict[str, Any]]:
    return _parse_delimited(file_path, delimiter=",", fmt="csv")


def _parse_tsv(file_path: str) -> Tuple[str, Dict[str, Any]]:
    return _parse_delimited(file_path, delimiter="\t", fmt="tsv")


def _parse_delimited(file_path: str, delimiter: str, fmt: str) -> Tuple[str, Dict[str, Any]]:
    rows: list = []
    try:
        with open(file_path, encoding="utf-8", errors="ignore", newline="") as fh:
            reader = csv.DictReader(fh, delimiter=delimiter)
            for row in reader:
                rows.append(row)
        content = "\n".join(
            "\t".join(f"{k}: {v}" for k, v in row.items()) for row in rows
        )
        return content, {"format": fmt, "row_count": len(rows)}
    except Exception as exc:
        logger.error("%s parsing failed for %s: %s", fmt.upper(), file_path, exc)
        return "", {"format": fmt, "error": str(exc)}


def _parse_json(file_path: str) -> Tuple[str, Dict[str, Any]]:
    try:
        with open(file_path, encoding="utf-8", errors="ignore") as fh:
            data = json.load(fh)
        content = json.dumps(data, indent=2, ensure_ascii=False)
        return content, {"format": "json"}
    except Exception as exc:
        logger.error("JSON parsing failed for %s: %s", file_path, exc)
        return "", {"format": "json", "error": str(exc)}


def _parse_jsonl(file_path: str) -> Tuple[str, Dict[str, Any]]:
    lines = []
    try:
        with open(file_path, encoding="utf-8", errors="ignore") as fh:
            for line in fh:
                line = line.strip()
                if line:
                    try:
                        obj = json.loads(line)
                        lines.append(json.dumps(obj, ensure_ascii=False))
                    except json.JSONDecodeError:
                        lines.append(line)
        return "\n".join(lines), {"format": "jsonl", "record_count": len(lines)}
    except Exception as exc:
        logger.error("JSONL parsing failed for %s: %s", file_path, exc)
        return "", {"format": "jsonl", "error": str(exc)}


def _parse_html(file_path: str) -> Tuple[str, Dict[str, Any]]:
    try:
        from bs4 import BeautifulSoup
        with open(file_path, encoding="utf-8", errors="ignore") as fh:
            soup = BeautifulSoup(fh.read(), "html.parser")
        # Remove script and style elements
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        content = soup.get_text(separator="\n", strip=True)
        title = soup.title.string if soup.title else ""
        return content, {"format": "html", "title": title}
    except ImportError:
        logger.warning("beautifulsoup4 not installed — falling back to plain-text for %s", file_path)
        return _parse_text(file_path)
    except Exception as exc:
        logger.error("HTML parsing failed for %s: %s", file_path, exc)
        return "", {"format": "html", "error": str(exc)}
